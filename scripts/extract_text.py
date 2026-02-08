#!/usr/bin/env python3
"""
00_inbox/ 内のバイナリファイル(.pptx, .xlsx, .xlsm, .pdf)からテキストを抽出し、
構造化 JSON として /tmp/extracted_texts.json に出力する。

GitHub Actions の analyze-inbox ワークフローから呼び出される。
ローカルでも単体実行可能。
"""

import json
import os
import sys
from pathlib import Path

INBOX_DIR = Path("00_inbox")
OUTPUT_PATH = Path("/tmp/extracted_texts.json")
MAX_XLSX_ROWS = 500


def extract_pptx(filepath: Path) -> dict:
    """PowerPoint ファイルからテキストを抽出する。"""
    from pptx import Presentation

    prs = Presentation(str(filepath))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(row_texts))
        if texts:
            slides.append({"slide_number": i, "content": texts})

    return {
        "type": "pptx",
        "slide_count": len(prs.slides),
        "slides": slides,
    }


def extract_xlsx(filepath: Path) -> dict:
    """Excel ファイルからテキストを抽出する。"""
    from openpyxl import load_workbook

    wb = load_workbook(str(filepath), read_only=True, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_data = []
        row_count = 0
        for row in ws.iter_rows(values_only=True):
            if row_count >= MAX_XLSX_ROWS:
                break
            cells = [str(cell) if cell is not None else "" for cell in row]
            if any(c for c in cells):
                rows_data.append(cells)
            row_count += 1
        if rows_data:
            sheets.append({
                "sheet_name": sheet_name,
                "row_count": len(rows_data),
                "truncated": row_count >= MAX_XLSX_ROWS,
                "rows": rows_data,
            })
    sheet_count = len(wb.sheetnames)
    wb.close()

    return {
        "type": "xlsx",
        "sheet_count": sheet_count,
        "sheets": sheets,
    }


def extract_pdf(filepath: Path) -> dict:
    """PDF ファイルからテキストを抽出する。"""
    from PyPDF2 import PdfReader

    reader = PdfReader(str(filepath))
    pages = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text and text.strip():
            pages.append({"page_number": i, "content": text.strip()})

    return {
        "type": "pdf",
        "page_count": len(reader.pages),
        "pages": pages,
    }


EXTRACTORS = {
    ".pptx": extract_pptx,
    ".xlsx": extract_xlsx,
    ".xlsm": extract_xlsx,
    ".pdf": extract_pdf,
}


def main():
    if not INBOX_DIR.exists():
        print(f"Warning: {INBOX_DIR} does not exist")
        sys.exit(0)

    files = [
        f for f in INBOX_DIR.rglob("*")
        if f.is_file() and f.name != ".gitkeep"
    ]

    if not files:
        print("No files found in 00_inbox/")
        sys.exit(0)

    results = []
    for filepath in files:
        print(f"Processing: {filepath}")
        ext = filepath.suffix.lower()
        file_info = {
            "filename": filepath.name,
            "path": str(filepath),
            "size_bytes": filepath.stat().st_size,
            "extension": ext,
        }

        extractor = EXTRACTORS.get(ext)
        if extractor:
            try:
                extracted = extractor(filepath)
                file_info["extracted"] = extracted
                file_info["status"] = "success"
            except Exception as e:
                file_info["status"] = "error"
                file_info["error"] = str(e)
                print(f"  Error extracting {filepath}: {e}", file=sys.stderr)
        else:
            file_info["status"] = "unsupported"
            file_info["note"] = f"No extractor for {ext}. File will be classified by filename."

        results.append(file_info)

    # Windows 互換: ローカル実行時は /tmp が無い場合カレントに出力
    output_path = OUTPUT_PATH
    if sys.platform == "win32" and not OUTPUT_PATH.parent.exists():
        output_path = Path("extracted_texts.json")

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(results, f, ensure_ascii=False, indent=2)

    print(f"\nExtracted {len(results)} file(s) -> {output_path}")


if __name__ == "__main__":
    main()
